import os
import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_full_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # Global Colors
    c_header_bg = RGBColor(30, 77, 123)       # Dark Blue
    c_number_bg = RGBColor(26, 155, 252)      # Sky Blue
    c_white = RGBColor(255, 255, 255)
    c_dark_gray = RGBColor(40, 40, 40)
    c_text_blue = RGBColor(20, 50, 90)
    c_teal_green = RGBColor(12, 128, 112)     # Green/Teal for Gap
    
    # Helper to add standard slide header
    def add_slide_header(slide, number, title, category):
        # Header Background
        h_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.0))
        h_bg.fill.solid()
        h_bg.fill.fore_color.rgb = c_header_bg
        h_bg.line.fill.background()
        
        # Number Box
        num_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(0.15), Inches(0.7), Inches(0.7))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = c_number_bg
        num_box.line.fill.background()
        tf_num = num_box.text_frame
        p_num = tf_num.paragraphs[0]
        p_num.alignment = PP_ALIGN.CENTER
        p_num.text = number
        p_num.font.name = "Arial"
        p_num.font.size = Pt(28)
        p_num.font.bold = True
        p_num.font.color.rgb = c_white
        
        # Header Title
        title_box = slide.shapes.add_textbox(Inches(1.3), Inches(0.15), Inches(9.0), Inches(0.7))
        tf_title = title_box.text_frame
        p_title = tf_title.paragraphs[0]
        p_title.text = title
        p_title.font.name = "Arial"
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = c_white
        
        # Header Category
        cat_box = slide.shapes.add_textbox(Inches(11.3), Inches(0.3), Inches(1.6), Inches(0.4))
        tf_cat = cat_box.text_frame
        p_cat = tf_cat.paragraphs[0]
        p_cat.alignment = PP_ALIGN.RIGHT
        p_cat.text = category
        p_cat.font.name = "Arial"
        p_cat.font.size = Pt(12)
        p_cat.font.bold = True
        p_cat.font.color.rgb = c_white

    # =========================================================================
    # SLIDE 04: LITERATURE REVIEW & ANALYSIS
    # =========================================================================
    slide04 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide04, "04", "LITERATURE REVIEW & ANALYSIS", "INTRODUCTION")
    
    # Subtitle text
    sub_box04 = slide04.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.333), Inches(0.4))
    tf_sub04 = sub_box04.text_frame
    p_sub04 = tf_sub04.paragraphs[0]
    p_sub04.text = "A systematic review across the project's domains — surveillance, on-device vision, model compression, IoT & sensing — with the gap each leaves open."
    p_sub04.font.name = "Arial"
    p_sub04.font.size = Pt(10.5)
    p_sub04.font.italic = True
    p_sub04.font.color.rgb = RGBColor(100, 100, 100)
    
    # Table data
    rows = 7
    cols = 4
    table_shape = slide04.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(12.333), Inches(3.9))
    table = table_shape.table
    
    # Set column widths
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(2.2)
    table.columns[2].width = Inches(3.9)
    table.columns[3].width = Inches(4.433)
    
    # Table headers
    headers = ["THEME", "KEY WORKS", "WHAT THE LITERATURE SHOWS", "GAP → OUR OPPORTUNITY"]
    for i, h_text in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h_text
        # styling
        cell.fill.solid()
        if i == 3:
            cell.fill.fore_color.rgb = c_teal_green
        else:
            cell.fill.fore_color.rgb = RGBColor(235, 242, 250)
            
        p = cell.text_frame.paragraphs[0]
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = c_white if i == 3 else c_text_blue
        p.alignment = PP_ALIGN.LEFT
        
    row_data = [
        ("Cloud / commercial cameras", "Ring, Nest; Al-Khedher [1]", "Continuous cloud streaming enables smart alerts.", "Privacy risk, high bandwidth, recurring fees."),
        ("Edge & privacy computing", "Shi et al. [2]; Tanwar et al. [3]", "Moving compute to the edge cuts latency & exposure.", "Most still offload inference; few run fully on-device."),
        ("Detection models", "YOLO [4], SSD [5] vs MobileNet [6][7]", "Accurate detectors exist; depthwise-separable conv is far lighter.", "Heavy nets exhaust 520 KB SRAM (1–3 FPS, brownouts); light nets target phones, not a ~100 KB MCU."),
        ("Model compression", "INT8 quant. [8], distillation [9]", "~4× smaller models with integer-only math.", "A real quantization gap; tiny models lose generalization."),
        ("TinyML on MCUs", "Warden [10]; TFLite-Micro [11]", "OS-free, low-latency inference on bare metal.", "Mostly research demos, not privacy-first product systems."),
        ("Sensing & connectivity", "PIR motion; MQTT [12]; UrbanSound8K [13]", "Cheap triggers; lightweight pub/sub; urban-sound classification.", "PIR alone = 85–95% false alarms; MQTT demos are LAN-only; audio rarely fused with vision.")
    ]
    
    for r_idx, data in enumerate(row_data):
        for c_idx, val in enumerate(data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            # Styling cell texts
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Arial"
            p.font.size = Pt(9.5)
            
            if c_idx == 0:
                p.font.bold = True
                p.font.color.rgb = c_text_blue
            elif c_idx == 1:
                p.font.color.rgb = RGBColor(90, 90, 90)
            elif c_idx == 2:
                p.font.color.rgb = c_dark_gray
            elif c_idx == 3:
                p.font.bold = True
                p.font.color.rgb = c_teal_green
                
            cell.fill.solid()
            # Alternating row background colors
            if r_idx % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(250, 252, 255)
            else:
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                
    # Synthesis/Research Gap Box
    gap_bg = slide04.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.6), Inches(12.333), Inches(1.4))
    gap_bg.fill.solid()
    gap_bg.fill.fore_color.rgb = RGBColor(25, 62, 107)  # Dark Blue-Gray
    gap_bg.line.fill.background()
    
    tf_gap = gap_bg.text_frame
    tf_gap.word_wrap = True
    tf_gap.margin_left = tf_gap.margin_right = Inches(0.2)
    tf_gap.margin_top = Inches(0.15)
    
    p_gap_header = tf_gap.paragraphs[0]
    p_gap_header.text = "RESEARCH GAP  —  the synthesis"
    p_gap_header.font.name = "Arial"
    p_gap_header.font.size = Pt(13)
    p_gap_header.font.bold = True
    p_gap_header.font.color.rgb = c_number_bg
    p_gap_header.space_after = Pt(4)
    
    p_gap_body = tf_gap.add_paragraph()
    p_gap_body.text = "No low-cost system combines (a) fully on-device human detection for privacy, (b) outbound-only cloud access that works on any network without port-forwarding, and (c) audio threat fusion — on a ~$10 microcontroller. EagleEye targets exactly this gap."
    p_gap_body.font.name = "Arial"
    p_gap_body.font.size = Pt(10.5)
    p_gap_body.font.color.rgb = c_white
    p_gap_body.space_after = Pt(2)

    # =========================================================================
    # SLIDE 05: FOUNDATIONS WE BUILD ON
    # =========================================================================
    slide05 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide05, "05", "FOUNDATIONS WE BUILD ON", "INTRODUCTION")
    
    # Subtitle text
    sub_box05 = slide05.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.333), Inches(0.4))
    tf_sub05 = sub_box05.text_frame
    p_sub05 = tf_sub05.paragraphs[0]
    p_sub05.text = "We don't just cite these techniques — we benchmarked and adapted them (18 model variants; see \"AI Vision\")."
    p_sub05.font.name = "Arial"
    p_sub05.font.size = Pt(10.5)
    p_sub05.font.bold = True
    p_sub05.font.color.rgb = c_dark_gray
    
    # 6 Blocks data with UPDATED TEXT inside
    blocks_data_s5 = [
        {
            "title": "Depthwise-separable convolutions [6][7]",
            "desc": "Factorize conv into depthwise + 1x1 pointwise -> ~8–10x fewer MACs.\n\n→ Enables our from-scratch on-device CNN (v7.16) based on MobileNet structure.",
            "fill": RGBColor(245, 250, 245), "border": RGBColor(185, 215, 185), "color": RGBColor(50, 120, 50)
        },
        {
            "title": "INT8 post-training quantization [8]",
            "desc": "32-bit -> 8-bit. ~75% smaller; integer-only math.\n\n→ Fits the ~100 KB tensor arena. Quantizes both weights & input pixels (Scale & Zero-Point).",
            "fill": RGBColor(245, 248, 252), "border": RGBColor(190, 210, 230), "color": RGBColor(50, 90, 140)
        },
        {
            "title": "Knowledge distillation [9]",
            "desc": "Transfer knowledge from a large teacher model to a small student model.\n\n→ Evaluated (v4.0, 91.77%); informed our capacity & resolution choices.",
            "fill": RGBColor(250, 245, 252), "border": RGBColor(215, 190, 225), "color": RGBColor(110, 60, 130)
        },
        {
            "title": "TFLite-Micro on bare metal [10][11]",
            "desc": "OS-free, low-latency inference directly on microcontrollers.\n\n→ Runs on ESP32-CAM. Accelerated by ESP-NN assembly kernels (3x speedup, 2.5s -> 0.8s).",
            "fill": RGBColor(253, 248, 243), "border": RGBColor(230, 205, 180), "color": RGBColor(140, 80, 40)
        },
        {
            "title": "Edge / privacy-by-design [2][3]",
            "desc": "Analyze locally on-chip; never stream raw video to the cloud.\n\n→ Transmit only confirmed human threat snapshots for alerts.",
            "fill": RGBColor(252, 245, 245), "border": RGBColor(225, 190, 190), "color": RGBColor(140, 50, 50)
        },
        {
            "title": "MQTT pub/sub [12] + UrbanSound8K [13]",
            "desc": "Lightweight pub/sub network messaging combined with audio classification.\n\n→ Outbound-only cloud alerts + triggers from our audio threat detector.",
            "fill": RGBColor(242, 249, 250), "border": RGBColor(185, 215, 220), "color": RGBColor(30, 110, 120)
        }
    ]
    
    # Draw the 6 boxes in 2 rows of 3 columns
    b_width = Inches(3.9)
    b_height = Inches(1.9)
    spacing_x = Inches(0.316)
    spacing_y = Inches(0.2)
    start_x = Inches(0.5)
    start_y = Inches(1.5)
    
    for idx, b in enumerate(blocks_data_s5):
        row = idx // 3
        col = idx % 3
        
        bx = start_x + col * (b_width + spacing_x)
        by = start_y + row * (b_height + spacing_y)
        
        # Draw background shape
        box = slide05.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, b_width, b_height)
        box.fill.solid()
        box.fill.fore_color.rgb = b["fill"]
        box.line.color.rgb = b["border"]
        box.line.width = Pt(1.5)
        
        # Add stylized circle inside box
        circle = slide05.shapes.add_shape(MSO_SHAPE.OVAL, bx + Inches(0.15), by + Inches(0.15), Inches(0.45), Inches(0.45))
        circle.fill.solid()
        circle.fill.fore_color.rgb = b["color"]
        circle.line.fill.background()
        
        # Text inside circle
        p_c = circle.text_frame.paragraphs[0]
        p_c.alignment = PP_ALIGN.CENTER
        p_c.text = ["🧩", "⚙️", "🎓", "🛠️", "🔒", "📡"][idx]
        p_c.font.size = Pt(14)
        
        # Add text box inside shape
        tb = slide05.shapes.add_textbox(bx + Inches(0.7), by + Inches(0.1), Inches(3.1), Inches(1.7))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = Inches(0)
        
        p_t = tf.paragraphs[0]
        p_t.text = b["title"]
        p_t.font.name = "Arial"
        p_t.font.size = Pt(11)
        p_t.font.bold = True
        p_t.font.color.rgb = c_text_blue
        p_t.space_after = Pt(4)
        
        p_d = tf.add_paragraph()
        p_d.text = b["desc"]
        p_d.font.name = "Arial"
        p_d.font.size = Pt(9.0)
        p_d.font.color.rgb = c_dark_gray
        
    # Bottom positioning block
    pos_bg = slide05.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.8), Inches(12.333), Inches(1.2))
    pos_bg.fill.solid()
    pos_bg.fill.fore_color.rgb = RGBColor(238, 247, 245)  # Very light mint
    pos_bg.line.color.rgb = RGBColor(170, 210, 200)
    pos_bg.line.width = Pt(1.5)
    
    tf_pos = pos_bg.text_frame
    tf_pos.word_wrap = True
    tf_pos.margin_left = tf_pos.margin_right = Inches(0.2)
    tf_pos.margin_top = Inches(0.15)
    
    p_pos_h = tf_pos.paragraphs[0]
    p_pos_h.text = "POSITIONING"
    p_pos_h.font.name = "Arial"
    p_pos_h.font.size = Pt(11)
    p_pos_h.font.bold = True
    p_pos_h.font.color.rgb = c_teal_green
    p_pos_h.space_after = Pt(4)
    
    p_pos_b = tf_pos.add_paragraph()
    p_pos_b.text = "EagleEye combines these proven building blocks into a single privacy-first, remotely-accessible edge device — and contributes a custom from-scratch depthwise-separable detector that holds >90% accuracy entirely on-device."
    p_pos_b.font.name = "Arial"
    p_pos_b.font.size = Pt(10.5)
    p_pos_b.font.color.rgb = c_dark_gray

    # =========================================================================
    # SLIDE 09: ESP32 CHIP ARCHITECTURE & MEMORY MAP
    # =========================================================================
    slide09 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide09, "09", "ESP32 CHIP ARCHITECTURE & MEMORY MAP", "FIRMWARE")
    
    blocks_data_s9 = [
        {
            "title": "FLASH (4 MB)",
            "left": Inches(0.5), "top": Inches(1.3), "width": Inches(4.5), "height": Inches(2.6),
            "fill": RGBColor(245, 247, 250), "border": RGBColor(180, 190, 210),
            "bullets": [
                ("Firmware code", 0, True),
                ("Model: 13 KB", 0, True),
                ("• architecture", 1, False),
                ("• weights", 1, False),
                ("ESP-NN library code", 0, True),
                ("• optimized C for conv, dense, pool", 1, False),
                ("• replaces generic TFLite math (3× faster)", 1, False)
            ]
        },
        {
            "title": "INTERNAL SRAM (~520 KB)",
            "left": Inches(5.5), "top": Inches(1.3), "width": Inches(7.3), "height": Inches(3.5),
            "fill": RGBColor(242, 249, 245), "border": RGBColor(170, 200, 185),
            "bullets": [
                ("Tensor Arena (126 KB)", 0, True),
                ("• conv outputs live here (reused every layer)", 1, False),
                ("snapshot_buf (27 KB)", 0, True),
                ("• 96×96 RGB888 image buffer", 1, False),
                ("TLS session (40 KB)", 0, True),
                ("• MQTT encryption overhead", 1, False),
                ("FreeRTOS + WiFi (~150 KB)", 0, True),
                ("• operating system & network stack", 1, False)
            ]
        },
        {
            "title": "PSRAM (4 MB)",
            "left": Inches(0.5), "top": Inches(4.1), "width": Inches(4.5), "height": Inches(1.4),
            "fill": RGBColor(249, 243, 251), "border": RGBColor(210, 180, 220),
            "bullets": [
                ("Camera buffer 1 (150 KB)", 0, False),
                ("Camera buffer 2 (150 KB)", 0, False),
                ("• dual-buffered DMA camera capture", 1, False)
            ]
        },
        {
            "title": "HARDWARE",
            "left": Inches(0.5), "top": Inches(5.7), "width": Inches(4.5), "height": Inches(1.5),
            "fill": RGBColor(253, 247, 241), "border": RGBColor(230, 195, 170),
            "bullets": [
                ("LEDC PWM → controls pan/tilt servos", 0, False),
                ("DMA → high-speed camera data transfer", 0, False),
                ("Radio → WiFi antenna signal", 0, False)
            ]
        },
        {
            "title": "CPU",
            "left": Inches(5.5), "top": Inches(5.0), "width": Inches(7.3), "height": Inches(2.2),
            "fill": RGBColor(242, 246, 252), "border": RGBColor(180, 205, 235),
            "bullets": [
                ("Core 0: servo stepper task (330Hz) | WiFi driver & stack", 0, True),
                ("Core 1: camera capture → resize → normalize → conv(ESP-NN) → pool(ESP-NN) → sepconv(ESP-NN) → pool(ESP-NN) → sepconv(ESP-NN) → pool → dense(ESP-NN) → decision", 0, True)
            ]
        }
    ]
    
    for b in blocks_data_s9:
        box = slide09.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, b["left"], b["top"], b["width"], b["height"])
        box.fill.solid()
        box.fill.fore_color.rgb = b["fill"]
        box.line.color.rgb = b["border"]
        box.line.width = Pt(1.5)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = Inches(0.15)
        
        p_title = tf.paragraphs[0]
        p_title.text = b["title"]
        p_title.font.name = "Arial"
        p_title.font.size = Pt(13)
        p_title.font.bold = True
        p_title.font.color.rgb = c_text_blue
        p_title.space_after = Pt(6)
        
        for text, level, bold in b["bullets"]:
            p = tf.add_paragraph()
            p.text = text
            p.level = level
            p.font.name = "Arial"
            p.font.size = Pt(10) if level == 1 else Pt(11)
            p.font.bold = bold
            p.font.color.rgb = c_dark_gray
            p.space_after = Pt(2)
            
    # Arrow for Slide 09
    arrow = slide09.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.1), Inches(2.3), Inches(0.35), Inches(0.25))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(100, 120, 150)
    arrow.line.fill.background()

    # =========================================================================
    # SLIDE 10: MODEL ACCELERATION & SPEED OPTIMIZATION
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide10, "10", "MODEL ACCELERATION & SPEED OPTIMIZATION", "MODEL")
    
    pillars_data = [
        {
            "title": "1. ESP-NN Library Optimization",
            "subtitle": "Hardware-Level Assembly Kernels",
            "left": Inches(0.6), "top": Inches(1.3), "width": Inches(3.8), "height": Inches(5.6),
            "fill": RGBColor(242, 246, 252), "border": RGBColor(180, 205, 235),
            "bullets": [
                ("Replaces standard TFLite math loops", 0, True),
                ("• Replaces generic C code with hand-optimized assembly routines from Espressif.", 1, False),
                ("Targets key pipeline layers", 0, True),
                ("• Optimizes Conv2D, MaxPool, DepthwiseConv, and Fully Connected operations.", 1, False),
                ("Xtensa LX6 SIMD Instructions", 0, True),
                ("• Executes multiple math operations simultaneously per CPU clock cycle.", 1, False),
                ("Massive Speedup Result", 0, True),
                ("• Reduces inference time from 2.5s down to 800ms (3× speed boost).", 1, False)
            ]
        },
        {
            "title": "2. INT8 Quantization",
            "subtitle": "Integer Math Conversion",
            "left": Inches(4.76), "top": Inches(1.3), "width": Inches(3.8), "height": Inches(5.6),
            "fill": RGBColor(242, 249, 245), "border": RGBColor(170, 200, 185),
            "bullets": [
                ("Converts weights & activations", 0, True),
                ("• Drops model float values from 32-bit floats to signed 8-bit integers.", 1, False),
                ("Bypasses FPU hardware limit", 0, True),
                ("• ESP32 lacks double-precision FPU; float operations require slow software emulation.", 1, False),
                ("Single-cycle ALU operation", 0, True),
                ("• Integer additions and multiplications run natively in single clock cycles.", 1, False),
                ("Saves critical Flash space", 0, True),
                ("• Quantized model weighs 13 KB instead of 52 KB (75% memory footprint reduction).", 1, False)
            ]
        },
        {
            "title": "3. Depthwise Separable Convs",
            "subtitle": "Architectural Layer Design",
            "left": Inches(8.92), "top": Inches(1.3), "width": Inches(3.8), "height": Inches(5.6),
            "fill": RGBColor(253, 247, 241), "border": RGBColor(230, 195, 170),
            "bullets": [
                ("Decomposes standard convolutions", 0, True),
                ("• Splits standard 2D convs into 2 distinct operations: Depthwise and Pointwise.", 1, False),
                ("Depthwise step (3x3 filter)", 0, True),
                ("• Filters input channels individually (applies spatial filter channel-by-channel).", 1, False),
                ("Pointwise step (1x1 projection)", 0, True),
                ("• Linearly combines outputs across all channels to build final features.", 1, False),
                ("9x Reduction in MAC Math", 0, True),
                ("• Drastically cuts parameters and CPU operations with minimal accuracy loss.", 1, False)
            ]
        }
    ]
    
    for p_data in pillars_data:
        box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, p_data["left"], p_data["top"], p_data["width"], p_data["height"])
        box.fill.solid()
        box.fill.fore_color.rgb = p_data["fill"]
        box.line.color.rgb = p_data["border"]
        box.line.width = Pt(1.5)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = Inches(0.15)
        
        # Add Column Title
        p_t = tf.paragraphs[0]
        p_t.text = p_data["title"]
        p_t.font.name = "Arial"
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = c_text_blue
        
        # Add Subtitle
        p_sub = tf.add_paragraph()
        p_sub.text = p_data["subtitle"]
        p_sub.font.name = "Arial"
        p_sub.font.size = Pt(10)
        p_sub.font.italic = True
        p_sub.font.color.rgb = RGBColor(100, 100, 100)
        p_sub.space_after = Pt(8)
        
        for text, level, bold in p_data["bullets"]:
            p = tf.add_paragraph()
            p.text = text
            p.level = level
            p.font.name = "Arial"
            p.font.size = Pt(9.5) if level == 1 else Pt(10.5)
            p.font.bold = bold
            p.font.color.rgb = c_dark_gray
            p.space_after = Pt(2)
            
    # Save Presentation
    output_ppt_path = "C:\\fyp-eagle-eye\\esp32_architecture.pptx"
    prs.save(output_ppt_path)
    print("Full deck saved successfully!")

if __name__ == "__main__":
    create_full_deck()
