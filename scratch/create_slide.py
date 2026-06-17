import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_slide():
    prs = Presentation()
    
    # Set slide dimensions to widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Add a blank slide
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Colors
    c_header_bg = RGBColor(30, 77, 123)       # Dark Blue
    c_number_bg = RGBColor(26, 155, 252)      # Sky Blue
    c_white = RGBColor(255, 255, 255)
    c_dark_gray = RGBColor(40, 40, 40)
    c_text_blue = RGBColor(20, 50, 90)
    
    # -------------------------------------------------------------
    # 1. Slide Header
    # -------------------------------------------------------------
    # Header Background
    header_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.0)
    )
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = c_header_bg
    header_bg.line.fill.background() # No border
    
    # Header Number Box (09)
    num_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(0.15), Inches(0.7), Inches(0.7)
    )
    num_box.fill.solid()
    num_box.fill.fore_color.rgb = c_number_bg
    num_box.line.fill.background()
    
    tf_num = num_box.text_frame
    tf_num.word_wrap = True
    p_num = tf_num.paragraphs[0]
    p_num.alignment = PP_ALIGN.CENTER
    p_num.text = "09"
    p_num.font.name = "Arial"
    p_num.font.size = Pt(28)
    p_num.font.bold = True
    p_num.font.color.rgb = c_white
    
    # Header Title
    title_box = slide.shapes.add_textbox(Inches(1.3), Inches(0.15), Inches(9.0), Inches(0.7))
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = "ESP32 CHIP ARCHITECTURE & MEMORY MAP"
    p_title.font.name = "Arial"
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = c_white
    
    # Header Category (FIRMWARE)
    cat_box = slide.shapes.add_textbox(Inches(11.3), Inches(0.3), Inches(1.6), Inches(0.4))
    tf_cat = cat_box.text_frame
    p_cat = tf_cat.paragraphs[0]
    p_cat.alignment = PP_ALIGN.RIGHT
    p_cat.text = "FIRMWARE"
    p_cat.font.name = "Arial"
    p_cat.font.size = Pt(12)
    p_cat.font.bold = True
    p_cat.font.color.rgb = c_white
    
    # -------------------------------------------------------------
    # 2. Block Configurations
    # -------------------------------------------------------------
    blocks_data = [
        {
            "title": "FLASH (4 MB)",
            "left": Inches(0.5), "top": Inches(1.3), "width": Inches(4.5), "height": Inches(2.6),
            "fill": RGBColor(245, 247, 250), "border": RGBColor(180, 190, 210),
            "bullets": [
                ("Firmware code", 0, True),
                ("Model: 13 KB", 0, True),
                ("• architecture", 1, False),
                ("• weights", 1, False),
                ("ESP-NN library code", 0, True),
                ("• optimized C for conv, dense, pool", 1, False),
                ("• replaces generic TFLite math (3× faster)", 1, False)
            ]
        },
        {
            "title": "INTERNAL SRAM (~520 KB)",
            "left": Inches(5.5), "top": Inches(1.3), "width": Inches(7.3), "height": Inches(3.5),
            "fill": RGBColor(242, 249, 245), "border": RGBColor(170, 200, 185),
            "bullets": [
                ("Tensor Arena (126 KB)", 0, True),
                ("• conv outputs live here (reused every layer)", 1, False),
                ("snapshot_buf (27 KB)", 0, True),
                ("• 96×96 RGB888 image buffer", 1, False),
                ("TLS session (40 KB)", 0, True),
                ("• MQTT encryption overhead", 1, False),
                ("FreeRTOS + WiFi (~150 KB)", 0, True),
                ("• operating system & network stack", 1, False)
            ]
        },
        {
            "title": "PSRAM (4 MB)",
            "left": Inches(0.5), "top": Inches(4.1), "width": Inches(4.5), "height": Inches(1.4),
            "fill": RGBColor(249, 243, 251), "border": RGBColor(210, 180, 220),
            "bullets": [
                ("Camera buffer 1 (150 KB)", 0, False),
                ("Camera buffer 2 (150 KB)", 0, False),
                ("• dual-buffered DMA camera capture", 1, False)
            ]
        },
        {
            "title": "HARDWARE",
            "left": Inches(0.5), "top": Inches(5.7), "width": Inches(4.5), "height": Inches(1.5),
            "fill": RGBColor(253, 247, 241), "border": RGBColor(230, 195, 170),
            "bullets": [
                ("LEDC PWM → controls pan/tilt servos", 0, False),
                ("DMA → high-speed camera data transfer", 0, False),
                ("Radio → WiFi antenna signal", 0, False)
            ]
        },
        {
            "title": "CPU",
            "left": Inches(5.5), "top": Inches(5.0), "width": Inches(7.3), "height": Inches(2.2),
            "fill": RGBColor(242, 246, 252), "border": RGBColor(180, 205, 235),
            "bullets": [
                ("Core 0: servo stepper task (330Hz) | WiFi driver & stack", 0, True),
                ("Core 1: camera capture → resize → normalize → conv(ESP-NN) → pool(ESP-NN) → sepconv(ESP-NN) → pool(ESP-NN) → sepconv(ESP-NN) → pool → dense(ESP-NN) → decision", 0, True)
            ]
        }
    ]
    
    # -------------------------------------------------------------
    # 3. Create Boxes and Content
    # -------------------------------------------------------------
    for b in blocks_data:
        # Create shape
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, b["left"], b["top"], b["width"], b["height"])
        box.fill.solid()
        box.fill.fore_color.rgb = b["fill"]
        box.line.color.rgb = b["border"]
        box.line.width = Pt(1.5)
        
        # Text Frame setup
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.15)
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_bottom = Inches(0.15)
        
        # Add Title Paragraph
        p_title = tf.paragraphs[0]
        p_title.text = b["title"]
        p_title.font.name = "Arial"
        p_title.font.size = Pt(13)
        p_title.font.bold = True
        p_title.font.color.rgb = c_text_blue
        p_title.space_after = Pt(6)
        
        # Add Bullet/Info Paragraphs
        for text, level, bold in b["bullets"]:
            p = tf.add_paragraph()
            p.text = text
            p.level = level
            p.font.name = "Arial"
            p.font.size = Pt(10) if level == 1 else Pt(11)
            p.font.bold = bold
            p.font.color.rgb = c_dark_gray
            p.space_after = Pt(2)
            
    # -------------------------------------------------------------
    # 4. Add Connection Arrow
    # -------------------------------------------------------------
    # Arrow from Flash to SRAM
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(5.1), Inches(2.3), Inches(0.35), Inches(0.25)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(100, 120, 150)
    arrow.line.fill.background()
    
    prs.save("esp32_architecture.pptx")
    print("PowerPoint presentation 'esp32_architecture.pptx' generated successfully!")

if __name__ == "__main__":
    create_slide()
