import os
import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import win32com.client

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # Global Colors
    c_header_bg = RGBColor(30, 77, 123)       # Dark Blue
    c_number_bg = RGBColor(26, 155, 252)      # Sky Blue
    c_white = RGBColor(255, 255, 255)
    c_dark_gray = RGBColor(40, 40, 40)
    c_text_blue = RGBColor(20, 50, 90)
    
    # =========================================================================
    # SLIDE 1: ESP32 CHIP ARCHITECTURE & MEMORY MAP
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    
    # Header Background
    header_bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.0))
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = c_header_bg
    header_bg.line.fill.background()
    
    # Number Box (09)
    num_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(0.15), Inches(0.7), Inches(0.7))
    num_box.fill.solid()
    num_box.fill.fore_color.rgb = c_number_bg
    num_box.line.fill.background()
    tf_num = num_box.text_frame
    p_num = tf_num.paragraphs[0]
    p_num.alignment = PP_ALIGN.CENTER
    p_num.text = "09"
    p_num.font.name = "Arial"
    p_num.font.size = Pt(28)
    p_num.font.bold = True
    p_num.font.color.rgb = c_white
    
    # Header Title
    title_box = slide1.shapes.add_textbox(Inches(1.3), Inches(0.15), Inches(9.0), Inches(0.7))
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = "ESP32 CHIP ARCHITECTURE & MEMORY MAP"
    p_title.font.name = "Arial"
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = c_white
    
    # Header Category (FIRMWARE)
    cat_box = slide1.shapes.add_textbox(Inches(11.3), Inches(0.3), Inches(1.6), Inches(0.4))
    tf_cat = cat_box.text_frame
    p_cat = tf_cat.paragraphs[0]
    p_cat.alignment = PP_ALIGN.RIGHT
    p_cat.text = "FIRMWARE"
    p_cat.font.name = "Arial"
    p_cat.font.size = Pt(12)
    p_cat.font.bold = True
    p_cat.font.color.rgb = c_white
    
    # Box configurations for Slide 1
    blocks_data_s1 = [
        {
            "title": "FLASH (4 MB)",
            "left": Inches(0.5), "top": Inches(1.3), "width": Inches(4.5), "height": Inches(2.6),
            "fill": RGBColor(245, 247, 250), "border": RGBColor(180, 190, 210),
            "bullets": [
                ("Firmware code", 0, True),
                ("Model: 13 KB", 0, True),
                ("• architecture", 1, False),
                ("• weights", 1, False),
                ("ESP-NN library code", 0, True),
                ("• optimized C for conv, dense, pool", 1, False),
                ("• replaces generic TFLite math (3× faster)", 1, False)
            ]
        },
        {
            "title": "INTERNAL SRAM (~520 KB)",
            "left": Inches(5.5), "top": Inches(1.3), "width": Inches(7.3), "height": Inches(3.5),
            "fill": RGBColor(242, 249, 245), "border": RGBColor(170, 200, 185),
            "bullets": [
                ("Tensor Arena (126 KB)", 0, True),
                ("• conv outputs live here (reused every layer)", 1, False),
                ("snapshot_buf (27 KB)", 0, True),
                ("• 96×96 RGB888 image buffer", 1, False),
                ("TLS session (40 KB)", 0, True),
                ("• MQTT encryption overhead", 1, False),
                ("FreeRTOS + WiFi (~150 KB)", 0, True),
                ("• operating system & network stack", 1, False)
            ]
        },
        {
            "title": "PSRAM (4 MB)",
            "left": Inches(0.5), "top": Inches(4.1), "width": Inches(4.5), "height": Inches(1.4),
            "fill": RGBColor(249, 243, 251), "border": RGBColor(210, 180, 220),
            "bullets": [
                ("Camera buffer 1 (150 KB)", 0, False),
                ("Camera buffer 2 (150 KB)", 0, False),
                ("• dual-buffered DMA camera capture", 1, False)
            ]
        },
        {
            "title": "HARDWARE",
            "left": Inches(0.5), "top": Inches(5.7), "width": Inches(4.5), "height": Inches(1.5),
            "fill": RGBColor(253, 247, 241), "border": RGBColor(230, 195, 170),
            "bullets": [
                ("LEDC PWM → controls pan/tilt servos", 0, False),
                ("DMA → high-speed camera data transfer", 0, False),
                ("Radio → WiFi antenna signal", 0, False)
            ]
        },
        {
            "title": "CPU",
            "left": Inches(5.5), "top": Inches(5.0), "width": Inches(7.3), "height": Inches(2.2),
            "fill": RGBColor(242, 246, 252), "border": RGBColor(180, 205, 235),
            "bullets": [
                ("Core 0: servo stepper task (330Hz) | WiFi driver & stack", 0, True),
                ("Core 1: camera capture → resize → normalize → conv(ESP-NN) → pool(ESP-NN) → sepconv(ESP-NN) → pool(ESP-NN) → sepconv(ESP-NN) → pool → dense(ESP-NN) → decision", 0, True)
            ]
        }
    ]
    
    for b in blocks_data_s1:
        box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, b["left"], b["top"], b["width"], b["height"])
        box.fill.solid()
        box.fill.fore_color.rgb = b["fill"]
        box.line.color.rgb = b["border"]
        box.line.width = Pt(1.5)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = Inches(0.15)
        
        p_title = tf.paragraphs[0]
        p_title.text = b["title"]
        p_title.font.name = "Arial"
        p_title.font.size = Pt(13)
        p_title.font.bold = True
        p_title.font.color.rgb = c_text_blue
        p_title.space_after = Pt(6)
        
        for text, level, bold in b["bullets"]:
            p = tf.add_paragraph()
            p.text = text
            p.level = level
            p.font.name = "Arial"
            p.font.size = Pt(10) if level == 1 else Pt(11)
            p.font.bold = bold
            p.font.color.rgb = c_dark_gray
            p.space_after = Pt(2)
            
    # Arrow for Slide 1
    arrow = slide1.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.1), Inches(2.3), Inches(0.35), Inches(0.25))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(100, 120, 150)
    arrow.line.fill.background()
    
    # =========================================================================
    # SLIDE 2: MODEL SPEED ACCELERATION & OPTIMIZATION
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    
    # Header Background
    header_bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.0))
    header_bg2.fill.solid()
    header_bg2.fill.fore_color.rgb = c_header_bg
    header_bg2.line.fill.background()
    
    # Number Box (10)
    num_box2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(0.15), Inches(0.7), Inches(0.7))
    num_box2.fill.solid()
    num_box2.fill.fore_color.rgb = c_number_bg
    num_box2.line.fill.background()
    tf_num2 = num_box2.text_frame
    p_num2 = tf_num2.paragraphs[0]
    p_num2.alignment = PP_ALIGN.CENTER
    p_num2.text = "10"
    p_num2.font.name = "Arial"
    p_num2.font.size = Pt(28)
    p_num2.font.bold = True
    p_num2.font.color.rgb = c_white
    
    # Header Title
    title_box2 = slide2.shapes.add_textbox(Inches(1.3), Inches(0.15), Inches(9.0), Inches(0.7))
    tf_title2 = title_box2.text_frame
    p_title2 = tf_title2.paragraphs[0]
    p_title2.text = "MODEL ACCELERATION & SPEED OPTIMIZATION"
    p_title2.font.name = "Arial"
    p_title2.font.size = Pt(22)
    p_title2.font.bold = True
    p_title2.font.color.rgb = c_white
    
    # Header Category (MODEL)
    cat_box2 = slide2.shapes.add_textbox(Inches(11.3), Inches(0.3), Inches(1.6), Inches(0.4))
    tf_cat2 = cat_box2.text_frame
    p_cat2 = tf_cat2.paragraphs[0]
    p_cat2.alignment = PP_ALIGN.RIGHT
    p_cat2.text = "MODEL"
    p_cat2.font.name = "Arial"
    p_cat2.font.size = Pt(12)
    p_cat2.font.bold = True
    p_cat2.font.color.rgb = c_white
    
    # 3 Columns for Slide 2 (The 3 Pillars of Speed)
    pillars_data = [
        {
            "title": "1. ESP-NN Library Optimization",
            "subtitle": "Hardware-Level Assembly Kernels",
            "left": Inches(0.6), "top": Inches(1.3), "width": Inches(3.8), "height": Inches(5.6),
            "fill": RGBColor(242, 246, 252), "border": RGBColor(180, 205, 235),
            "bullets": [
                ("Replaces standard TFLite math loops", 0, True),
                ("• Replaces generic C code with hand-optimized assembly routines from Espressif.", 1, False),
                ("Targets key pipeline layers", 0, True),
                ("• Optimizes Conv2D, MaxPool, DepthwiseConv, and Fully Connected operations.", 1, False),
                ("Xtensa LX6 SIMD Instructions", 0, True),
                ("• Executes multiple math operations simultaneously per CPU clock cycle.", 1, False),
                ("Massive Speedup Result", 0, True),
                ("• Reduces inference time from 2.5s down to 800ms (3× speed boost).", 1, False)
            ]
        },
        {
            "title": "2. INT8 Quantization",
            "subtitle": "Integer Math Conversion",
            "left": Inches(4.76), "top": Inches(1.3), "width": Inches(3.8), "height": Inches(5.6),
            "fill": RGBColor(242, 249, 245), "border": RGBColor(170, 200, 185),
            "bullets": [
                ("Converts weights & activations", 0, True),
                ("• Drops model float values from 32-bit floats to signed 8-bit integers.", 1, False),
                ("Bypasses FPU hardware limit", 0, True),
                ("• ESP32 lacks double-precision FPU; float operations require slow software emulation.", 1, False),
                ("Single-cycle ALU operation", 0, True),
                ("• Integer additions and multiplications run natively in single clock cycles.", 1, False),
                ("Saves critical Flash space", 0, True),
                ("• Quantized model weighs 13 KB instead of 52 KB (75% memory footprint reduction).", 1, False)
            ]
        },
        {
            "title": "3. Depthwise Separable Convs",
            "subtitle": "Architectural Layer Design",
            "left": Inches(8.92), "top": Inches(1.3), "width": Inches(3.8), "height": Inches(5.6),
            "fill": RGBColor(253, 247, 241), "border": RGBColor(230, 195, 170),
            "bullets": [
                ("Decomposes standard convolutions", 0, True),
                ("• Splits standard 2D convs into 2 distinct operations: Depthwise and Pointwise.", 1, False),
                ("Depthwise step (3x3 filter)", 0, True),
                ("• Filters input channels individually (applies spatial filter channel-by-channel).", 1, False),
                ("Pointwise step (1x1 projection)", 0, True),
                ("• Linearly combines outputs across all channels to build final features.", 1, False),
                ("9x Reduction in MAC Math", 0, True),
                ("• Drastically cuts parameters and CPU operations with minimal accuracy loss.", 1, False)
            ]
        }
    ]
    
    for p_data in pillars_data:
        box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, p_data["left"], p_data["top"], p_data["width"], p_data["height"])
        box.fill.solid()
        box.fill.fore_color.rgb = p_data["fill"]
        box.line.color.rgb = p_data["border"]
        box.line.width = Pt(1.5)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = Inches(0.15)
        
        # Add Column Title
        p_t = tf.paragraphs[0]
        p_t.text = p_data["title"]
        p_t.font.name = "Arial"
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = c_text_blue
        
        # Add Subtitle
        p_sub = tf.add_paragraph()
        p_sub.text = p_data["subtitle"]
        p_sub.font.name = "Arial"
        p_sub.font.size = Pt(10)
        p_sub.font.italic = True
        p_sub.font.color.rgb = RGBColor(100, 100, 100)
        p_sub.space_after = Pt(8)
        
        for text, level, bold in p_data["bullets"]:
            p = tf.add_paragraph()
            p.text = text
            p.level = level
            p.font.name = "Arial"
            p.font.size = Pt(9.5) if level == 1 else Pt(10.5)
            p.font.bold = bold
            p.font.color.rgb = c_dark_gray
            p.space_after = Pt(2)
            
    # Save Presentation
    output_ppt_path = "C:\\fyp-eagle-eye\\esp32_architecture.pptx"
    prs.save(output_ppt_path)
    print("Presentation saved successfully at:", output_ppt_path)

def export_slides_to_png():
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    ppt_path = "C:\\fyp-eagle-eye\\esp32_architecture.pptx"
    try:
        ppt = powerpoint.Presentations.Open(os.path.abspath(ppt_path), ReadOnly=True, WithWindow=False)
        # Export Slide 1
        ppt.Slides[1].Export("C:\\fyp-eagle-eye\\esp32_architecture_ppt_preview.png", "PNG")
        # Export Slide 2
        ppt.Slides[2].Export("C:\\fyp-eagle-eye\\esp32_model_speed_preview.png", "PNG")
        ppt.Close()
        print("Slides exported successfully!")
    except Exception as e:
        print("Error exporting slides:", e)
    finally:
        powerpoint.Quit()

if __name__ == "__main__":
    create_presentation()
    export_slides_to_png()
