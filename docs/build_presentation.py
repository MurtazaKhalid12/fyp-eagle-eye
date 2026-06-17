#!/usr/bin/env python3
"""
Build EagleEye_Presentation.pptx from PRESENTATION_SLIDES_CONTENT.md.

Single source of truth = the markdown plan; re-run this whenever the plan
changes. Uses the EagleEye brand palette from docs/POSTER_BRIEF.md.

    python docs/build_presentation.py
"""
from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

DOCS = Path(__file__).resolve().parent
SRC = DOCS / "PRESENTATION_SLIDES_CONTENT.md"
OUT = DOCS / "EagleEye_Presentation.pptx"

# EagleEye brand palette
DEEP_BLUE = RGBColor(0x1F, 0x4E, 0x79)
ACCENT = RGBColor(0x21, 0x96, 0xF3)
INK = RGBColor(0x1A, 0x1A, 0x1A)
SLATE = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PANEL = RGBColor(0xF4, 0xF6, 0xF8)
FONT = "Calibri"

INLINE = re.compile(r"(\*\*.+?\*\*|\*[^*]+\*|`[^`]+`)")


# --------------------------------------------------------------------------- #
#  Markdown parsing
# --------------------------------------------------------------------------- #
def parse_slides(text: str):
    marker = "# Presentation Slide Deck Content"
    text = text[text.index(marker) + len(marker):]
    parts = re.split(r"^## (Slide \d+:.*)$", text, flags=re.M)
    slides = []
    for i in range(1, len(parts), 2):
        header = parts[i].strip()
        title = header.split(":", 1)[1].strip() if ":" in header else header
        slides.append((title, parse_body(parts[i + 1])))
    return slides


def parse_table(buf):
    rows = []
    for row in buf:
        cells = [c.strip() for c in row.strip().strip("|").split("|")]
        if all(re.fullmatch(r":?-+:?", c) for c in cells):
            continue  # separator row
        rows.append(cells)
    return rows


def parse_body(body: str):
    blocks, table_buf = [], []

    def flush():
        if table_buf:
            blocks.append(("table", parse_table(table_buf)))
            table_buf.clear()

    for raw in body.split("\n"):
        line = raw.rstrip()
        s = line.strip()
        if not s:
            flush()
            continue
        if set(s) == {"-"} and len(s) >= 3:  # '---' separator
            flush()
            continue
        if s.startswith("|"):
            table_buf.append(s)
            continue
        flush()
        if re.match(r"^\*\*\[(Visual|Diagram).*\]\*\*$", s):
            blocks.append(("visual", s.strip("*")))
            continue
        m = re.match(r"^(\s*)(?:[-*]|\d+\.)\s+(.*)$", raw)
        if m:
            level = min(len(m.group(1)) // 4, 3)
            blocks.append(("bullet", level, m.group(2)))
            continue
        blocks.append(("para", s))
    flush()
    return blocks


# --------------------------------------------------------------------------- #
#  Rendering helpers
# --------------------------------------------------------------------------- #
def add_runs(p, s, size, color=INK, bold=False, italic=False):
    for tok in INLINE.split(s):
        if not tok:
            continue
        b, i, t = bold, italic, tok
        if tok.startswith("**") and tok.endswith("**"):
            b, t = True, tok[2:-2]
        elif tok.startswith("*") and tok.endswith("*"):
            i, t = True, tok[1:-1]
        elif tok.startswith("`") and tok.endswith("`"):
            t = tok[1:-1]
        r = p.add_run()
        r.text = t
        r.font.size = size
        r.font.bold = b
        r.font.italic = i
        r.font.color.rgb = color
        r.font.name = FONT


def no_bullet(p):
    pPr = p._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        e = pPr.find(qn(tag))
        if e is not None:
            pPr.remove(e)
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def bullet_char(p, color=ACCENT):
    pPr = p._p.get_or_add_pPr()
    buFont = pPr.makeelement(qn("a:buFont"), {"typeface": FONT})
    buChar = pPr.makeelement(qn("a:buChar"), {"char": "•"})
    pPr.append(buFont)
    pPr.append(buChar)


def title_bar(slide, title):
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.30), Inches(12.2), Inches(0.95))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    add_runs(p, title, Pt(28), color=DEEP_BLUE, bold=True)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.28),
                                 Inches(2.0), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False


def render_table(slide, rows, top):
    nrows, ncols = len(rows), len(rows[0])
    gt = slide.shapes.add_table(nrows, ncols, Inches(0.7), top,
                                Inches(11.9), Inches(0.4 * nrows)).table
    for r in range(nrows):
        for c in range(ncols):
            cell = gt.cell(r, c)
            cell.margin_top = Pt(2)
            cell.margin_bottom = Pt(2)
            para = cell.text_frame.paragraphs[0]
            head = r == 0
            add_runs(para, rows[r][c], Pt(12),
                     color=WHITE if head else INK, bold=head)
            cell.fill.solid()
            cell.fill.fore_color.rgb = DEEP_BLUE if head else (PANEL if r % 2 else WHITE)


# --------------------------------------------------------------------------- #
#  Slide builders
# --------------------------------------------------------------------------- #
def build_title_slide(slide, blocks):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DEEP_BLUE
    kv = {}
    for b in blocks:
        if b[0] == "bullet":
            m = re.match(r"\*\*(.+?):\*\*\s*(.*)", b[2])
            if m:
                kv[m.group(1).strip().lower()] = m.group(2).strip()
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(11.3), Inches(3.3))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_runs(p, kv.get("title", "EagleEye"), Pt(40), color=WHITE, bold=True)
    no_bullet(p)
    for key, size in (("subtitle", 22), ("presenters", 18),
                      ("supervisor", 18), ("date", 16)):
        if kv.get(key):
            q = tf.add_paragraph()
            q.alignment = PP_ALIGN.CENTER
            q.space_before = Pt(14 if key == "subtitle" else 6)
            add_runs(q, kv[key], Pt(size),
                     color=RGBColor(0xCF, 0xE2, 0xF3), italic=(key == "subtitle"))
            no_bullet(q)


def build_content_slide(slide, title, blocks):
    title_bar(slide, title)
    notes = []
    text_blocks = [b for b in blocks if b[0] != "table"]
    tables = [b for b in blocks if b[0] == "table"]

    n = sum(1 for b in text_blocks if b[0] in ("bullet", "para"))
    base = Pt(17) if n <= 6 else (Pt(15) if n <= 9 else Pt(13))

    body_h = 3.4 if tables else 5.6
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.9), Inches(body_h))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for b in text_blocks:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(5)
        if b[0] == "para":
            add_runs(p, b[1], base, color=INK)
            no_bullet(p)
        elif b[0] == "visual":
            add_runs(p, "▦ " + b[1], Pt(12), color=SLATE, italic=True)
            no_bullet(p)
            notes.append(b[1])
        elif b[0] == "bullet":
            level, txt = b[1], b[2]
            p.level = level
            add_runs(p, txt, base if level == 0 else Pt(base.pt - 2),
                     color=INK if level == 0 else SLATE)
            bullet_char(p)

    for t in tables:
        render_table(slide, t[1], Inches(5.0))

    if notes:
        slide.notes_slide.notes_text_frame.text = "Visuals to add:\n- " + "\n- ".join(notes)


# --------------------------------------------------------------------------- #
def main():
    slides = parse_slides(SRC.read_text(encoding="utf-8"))
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for idx, (title, blocks) in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        if idx == 0 or title.lower().startswith("title slide"):
            build_title_slide(slide, blocks)
        else:
            build_content_slide(slide, title, blocks)

    prs.save(OUT)
    print(f"Wrote {OUT}  ({len(slides)} slides)")


if __name__ == "__main__":
    main()
