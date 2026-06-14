# -*- coding: utf-8 -*-
"""
Build the EagleEye A1 portrait academic poster as an editable PowerPoint (.pptx).
Clean-academic style: ITU blue/teal, embedded logo + hero + flowchart + repo QR,
native (editable) cloud-architecture diagram, metric tiles, dashed slots for missing assets.
"""
import os
import qrcode
from PIL import Image
from pptx import Presentation
from pptx.util import Mm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

DOCS = r"C:\fyp-eagle-eye\docs"
ASSETS = os.path.join(DOCS, "assets")
OUT = os.path.join(DOCS, "EagleEye_Poster_A1.pptx")

# ---------- palette ----------
BLUE  = RGBColor(0x1F, 0x4E, 0x79)
TEAL  = RGBColor(0x21, 0x96, 0xF3)
INK   = RGBColor(0x1A, 0x1A, 0x1A)
SLATE = RGBColor(0x55, 0x55, 0x55)
PANEL = RGBColor(0xF4, 0xF6, 0xF8)
RULE  = RGBColor(0xD0, 0xD5, 0xDB)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED   = RGBColor(0xE5, 0x39, 0x35)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LBLUE = RGBColor(0xEA, 0xF4, 0xFF)
FONT  = "Calibri"

prs = Presentation()
prs.slide_width = Mm(594)
prs.slide_height = Mm(841)
slide = prs.slides.add_slide(prs.slide_layouts[6])


# ---------- helpers ----------
def rect(x, y, w, h, fill=None, line=None, line_w=0.75, rounded=False, dash=False):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(st, Mm(x), Mm(y), Mm(w), Mm(h))
    if rounded:
        try: sp.adjustments[0] = 0.06
        except Exception: pass
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    if dash:
        ln = sp.line._get_or_add_ln()
        for el in ln.findall(qn('a:prstDash')): ln.remove(el)
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    sp.shadow.inherit = False
    return sp


def autoshape(shape_type, x, y, w, h, fill=None, line=None, line_w=0.75):
    sp = slide.shapes.add_shape(shape_type, Mm(x), Mm(y), Mm(w), Mm(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def _anchor(v):
    return {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}[v]


def text(x, y, w, h, paras, anchor="top", wrap=True):
    tb = slide.shapes.add_textbox(Mm(x), Mm(y), Mm(w), Mm(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = _anchor(anchor)
    tf.margin_left = Mm(1.5); tf.margin_right = Mm(1.5)
    tf.margin_top = Mm(0.8); tf.margin_bottom = Mm(0.8)
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get("align", PP_ALIGN.LEFT)
        if para.get("before") is not None: p.space_before = Pt(para["before"])
        if para.get("after") is not None: p.space_after = Pt(para["after"])
        if para.get("leading") is not None: p.line_spacing = para["leading"]
        for (txt, fmt) in para["runs"]:
            r = p.add_run(); r.text = txt
            f = r.font
            f.name = fmt.get("name", FONT)
            f.size = Pt(fmt.get("size", 18))
            f.bold = fmt.get("bold", False)
            f.italic = fmt.get("italic", False)
            f.color.rgb = fmt.get("color", INK)
    return tb


def R(txt, size=18, color=INK, bold=False, italic=False, name=FONT):
    return (txt, {"size": size, "color": color, "bold": bold, "italic": italic, "name": name})


def place_image(path, x, y, w, h, halign="center", valign="middle"):
    iw, ih = Image.open(path).size
    box_ar, img_ar = w / h, iw / ih
    if img_ar > box_ar: nw, nh = w, w / img_ar
    else: nh, nw = h, h * img_ar
    nx = x + (w - nw) * (0.5 if halign == "center" else 0.0)
    ny = y + (h - nh) * (0.5 if valign == "middle" else 0.0)
    return slide.shapes.add_picture(path, Mm(nx), Mm(ny), Mm(nw), Mm(nh))


def section(x, y, w, num, title, h=15):
    rect(x, y, w, h, fill=BLUE, rounded=True)
    text(x + 2, y, w - 4, h,
         [{"runs": [R(num + "  ", size=23, color=WHITE, bold=True),
                    R(title, size=22, color=WHITE, bold=True)], "align": PP_ALIGN.LEFT}],
         anchor="middle")
    return y + h + 5


def slot(x, y, w, h, label, sub="", fn=""):
    rect(x, y, w, h, fill=LBLUE, line=TEAL, line_w=2.0, rounded=True, dash=True)
    paras = [{"runs": [R(label, size=16, color=BLUE, bold=True)], "align": PP_ALIGN.CENTER, "after": 3}]
    if sub: paras.append({"runs": [R(sub, size=12, color=SLATE, italic=True)], "align": PP_ALIGN.CENTER, "after": 2})
    if fn: paras.append({"runs": [R(fn, size=10.5, color=SLATE, name="Consolas")], "align": PP_ALIGN.CENTER})
    text(x + 2, y, w - 4, h, paras, anchor="middle")


def img_card(path, x, y, w, h, caption):
    rect(x, y, w, h, fill=PANEL, line=RULE, line_w=1.0, rounded=True)
    if os.path.exists(path):
        place_image(path, x + 3, y + 3, w - 6, h - 14)
    text(x, y + h - 12, w, 11, [{"runs": [R(caption, size=10.5, color=SLATE, italic=True)], "align": PP_ALIGN.CENTER}])


P = lambda n: os.path.join(ASSETS, n)

# ================= QR =================
qr_repo = P("qr_repo.png")
qrcode.make("https://github.com/muhammadAB123/fyp-eagle-eye").save(qr_repo)

# ================= HEADER =================
rect(0, 0, 594, 122, fill=BLUE)
rect(14, 20, 152, 82, fill=WHITE, rounded=True)
if os.path.exists(P("itu_logo.png")):
    place_image(P("itu_logo.png"), 20, 27, 140, 68)
text(178, 12, 288, 102,
     [{"runs": [R("EagleEye", size=70, color=WHITE, bold=True)], "after": 2},
      {"runs": [R("Intelligent Edge AI Surveillance System", size=28, color=LBLUE, bold=True)], "after": 5},
      {"runs": [R("On-device AI   ·   Cloud-connected   ·   Watch your space from anywhere", size=16, color=WHITE, italic=True)]}],
     anchor="middle")
slot(470, 20, 110, 82, "SparkUp / ORIC", "logo", "assets/sparkup_logo.png")

rect(0, 122, 594, 42, fill=TEAL)
text(8, 122, 578, 42,
     [{"runs": [R("Murtaza Khalid (BSCE22004)      ·      Huzaifa Khan (BSCE22025)      ·      Haseeb Ahmed (BSCE22048)", size=17, color=WHITE, bold=True)], "align": PP_ALIGN.CENTER, "after": 2},
      {"runs": [R("Supervisor: Dr. Rehan Hafiz   ·   Co-Supervisor: Dr. Rehan Ahmed   ·   Dept. of Computer & Software Engineering, Information Technology University (ITU), Lahore", size=13.5, color=WHITE)], "align": PP_ALIGN.CENTER}],
     anchor="middle")

# ================= COLUMNS =================
LX, LW = 14, 230
RX, RW = 256, 324

# ---------------- LEFT ----------------
# (1) Problem
cy = section(LX, 176, LW, "①", "THE PROBLEM")
text(LX, cy, LW, 70, [
    {"runs": [R("Security cameras force a bad trade-off:", size=17, color=INK, bold=True)], "after": 6},
    {"runs": [R("Cloud cams", size=16, color=BLUE, bold=True), R(" stream your footage to a company's servers — privacy risk + monthly fees.", size=16)], "after": 6, "leading": 1.02},
    {"runs": [R("LAN / DIY cams", size=16, color=BLUE, bold=True), R(" work only on the same Wi-Fi — typed IP, a PC left on, risky port-forwarding.", size=16)], "after": 6, "leading": 1.02},
    {"runs": [R("Most ", size=16), R("record everything", size=16, color=BLUE, bold=True), R(" — real events buried in hours of empty video.", size=16)], "leading": 1.02},
])
rect(LX, 268, LW, 18, fill=LBLUE, line=TEAL, line_w=1.25, rounded=True)
text(LX, 268, LW, 18, [{"runs": [R("Goal:  ", size=15, color=TEAL, bold=True), R("see real intrusions — anywhere, privately, no subscription.", size=15, color=INK, italic=True)], "align": PP_ALIGN.CENTER}], anchor="middle")

# (2) Solution
cy = section(LX, 296, LW, "②", "OUR SOLUTION")
text(LX, cy, LW, 86, [
    {"runs": [R("EagleEye", size=16, color=BLUE, bold=True), R(" runs a neural network ", size=16), R("on the camera itself", size=16, bold=True), R(". A low-cost ESP32-CAM decides on-device whether a human is present — images are analysed locally, never streamed to the cloud.", size=16)], "after": 6, "leading": 1.02},
    {"runs": [R("On an intruder it pushes a ", size=16), R("photo alert to your phone", size=16, bold=True), R("; you can open ", size=16), R("live video", size=16, bold=True), R(" and ", size=16), R("pan the camera", size=16, bold=True), R(" — from any network (Wi-Fi or 4G/LTE).", size=16)], "after": 6, "leading": 1.02},
    {"runs": [R("The camera ", size=16), R("reaches OUT", size=16, color=TEAL, bold=True), R(" to the cloud — no static IP, no port-forwarding, no on-site PC. Encrypted end-to-end; runs on free tiers.", size=16)], "leading": 1.02},
])
img_card(P("hardware_setup.png"), LX, 392, LW, 178, "EagleEye node — ESP32-CAM in a 3D-printed mount  ·  swap for a real photo (assets/prototype_photo.jpg)")

# (3) Key features
cy = section(LX, 582, LW, "③", "KEY FEATURES")
feats = [
    ("On-device human detection", "TinyML; images never leave the device"),
    ("Instant intruder alerts", "captured photo pushed to your phone"),
    ("On-demand remote live video", "stream only when you ask"),
    ("App-controlled camera pan", "servo aim from the app"),
    ("Arm / disarm + true online status", "real online/offline, not guesswork"),
    ("Encrypted & private", "TLS everywhere, per-device credentials"),
    ("~$10 hardware  ·  $0/month cloud", "entire stack on free tiers"),
]
fp = [{"runs": [R("✓  ", size=17, color=GREEN, bold=True), R(a + "  ", size=16, color=INK, bold=True), R("— " + b, size=14, color=SLATE)], "after": 7, "leading": 1.02} for a, b in feats]
text(LX, cy, LW, 110, fp)

# (impact / SDG callout — fills the bottom of the left column)
rect(LX, 700, LW, 80, fill=LBLUE, line=TEAL, line_w=1.5, rounded=True)
text(LX + 3, 703, LW - 6, 74, [
    {"runs": [R("IMPACT & SUSTAINABILITY", size=15, color=BLUE, bold=True)], "after": 4},
    {"runs": [R("Affordable AI security", size=14, bold=True), R(" — ~$10 vs $200+ systems democratizes safety for low-income homes  (SDG 9 · 11).", size=14)], "after": 4, "leading": 1.02},
    {"runs": [R("Privacy by design", size=14, bold=True), R(" — only confirmed threats ever leave the device  (SDG 16).", size=14)], "leading": 1.02},
])

# ---------------- RIGHT ----------------
# (4) Architecture
cy = section(RX, 176, RW, "④", "SYSTEM ARCHITECTURE")
ax = RX
text(ax, cy, 90, 9, [{"runs": [R("SITE A", size=14, color=BLUE, bold=True)], "align": PP_ALIGN.CENTER}])
text(ax + 116, cy, 122, 9, [{"runs": [R("CLOUD", size=14, color=BLUE, bold=True)], "align": PP_ALIGN.CENTER}])
text(ax + 260, cy, 64, 9, [{"runs": [R("ANYWHERE", size=14, color=BLUE, bold=True)], "align": PP_ALIGN.CENTER}])
zy = cy + 11
rect(ax, zy, 90, 92, fill=PANEL, line=BLUE, line_w=1.75, rounded=True)
text(ax, zy, 90, 92, [
    {"runs": [R("ESP32-CAM", size=16, color=BLUE, bold=True)], "align": PP_ALIGN.CENTER, "after": 4},
    {"runs": [R("On-device AI", size=13)], "align": PP_ALIGN.CENTER, "after": 2},
    {"runs": [R("Servo pan", size=13)], "align": PP_ALIGN.CENTER, "after": 2},
    {"runs": [R("Captures JPEG", size=13)], "align": PP_ALIGN.CENTER}], anchor="middle")
autoshape(MSO_SHAPE.RIGHT_ARROW, ax + 92, zy + 38, 22, 16, fill=TEAL)
text(ax + 88, zy + 24, 30, 9, [{"runs": [R("TLS", size=11, color=SLATE, bold=True)], "align": PP_ALIGN.CENTER}])
cbx, cbw = ax + 116, 122
for j, lab in enumerate(["HiveMQ — MQTT broker", "Cloudinary  ·  Firebase", "Deno relay (live video)"]):
    by = zy + j * 32
    rect(cbx, by, cbw, 26, fill=BLUE, rounded=True)
    text(cbx, by, cbw, 26, [{"runs": [R(lab, size=13.5, color=WHITE, bold=True)], "align": PP_ALIGN.CENTER}], anchor="middle")
autoshape(MSO_SHAPE.RIGHT_ARROW, cbx + cbw, zy + 38, 22, 16, fill=TEAL)
text(cbx + cbw - 2, zy + 24, 30, 9, [{"runs": [R("wss", size=11, color=SLATE, bold=True)], "align": PP_ALIGN.CENTER}])
phx = cbx + cbw + 26
rect(phx, zy, 64, 92, fill=PANEL, line=BLUE, line_w=1.75, rounded=True)
text(phx, zy, 64, 92, [
    {"runs": [R("Phone app", size=15, color=BLUE, bold=True)], "align": PP_ALIGN.CENTER, "after": 2},
    {"runs": [R("React Native", size=11.5, color=SLATE)], "align": PP_ALIGN.CENTER, "after": 4},
    {"runs": [R("Alerts", size=13)], "align": PP_ALIGN.CENTER},
    {"runs": [R("Live view", size=13)], "align": PP_ALIGN.CENTER},
    {"runs": [R("Pan control", size=13)], "align": PP_ALIGN.CENTER}], anchor="middle")
text(ax, zy + 98, RW, 24, [
    {"runs": [R("The camera reaches OUT to the cloud; the phone reaches the same cloud — they meet in the middle.", size=14, color=INK, bold=True)], "align": PP_ALIGN.CENTER, "after": 2},
    {"runs": [R("Works on any network  ·  no static IP  ·  no port-forwarding  ·  no on-site PC  ·  end-to-end TLS", size=12.5, color=SLATE, italic=True)], "align": PP_ALIGN.CENTER}])

# (5) How it works
cy = section(RX, 348, RW, "⑤", "HOW IT WORKS")
stages = ["Capture frame", "AI: human?", "Snap JPEG", "Upload + Alert", "Phone push"]
sbw, sgap = 56, 8
sx = RX
for i, s in enumerate(stages):
    rect(sx, cy, sbw, 30, fill=LBLUE, line=TEAL, line_w=1.25, rounded=True)
    text(sx, cy, sbw, 30, [{"runs": [R(s, size=12.5, color=BLUE, bold=True)], "align": PP_ALIGN.CENTER}], anchor="middle")
    if i < len(stages) - 1:
        autoshape(MSO_SHAPE.RIGHT_ARROW, sx + sbw, cy + 10, sgap, 10, fill=TEAL)
    sx += sbw + sgap
text(RX, cy + 34, RW, 12, [{"runs": [R("On request → ", size=14, color=INK, bold=True), R("live video via cloud relay", size=14, color=TEAL, bold=True), R("      |      No human → keep watching  (only confirmed humans trigger an alert).", size=14, color=SLATE)]}])

# (6) Results & tech
cy = section(RX, 416, RW, "⑥", "RESULTS & TECH STACK")
tiles = [
    ("90.8%", "detection accuracy (test; 87.5% human recall)"),
    ("<1 s", "on-device inference (~872 ms, ESP32-CAM)"),
    ("<3 s", "detection → phone alert (end-to-end)"),
    ("18", "model versions (v1.0 → v7.17)"),
    ("~$10", "device vs $200+ · $0/month cloud"),
    ("100%", "private — frames stay on device"),
]
tw = (RW - 8 * 2) / 3
th = 52
for i, (n, l) in enumerate(tiles):
    tx = RX + (i % 3) * (tw + 8)
    ty = cy + (i // 3) * (th + 8)
    rect(tx, ty, tw, th, fill=PANEL, line=RULE, line_w=1.0, rounded=True)
    text(tx, ty + 3, tw, th - 4, [
        {"runs": [R(n, size=36, color=BLUE, bold=True)], "align": PP_ALIGN.CENTER, "after": 1},
        {"runs": [R(l, size=11.5, color=SLATE)], "align": PP_ALIGN.CENTER, "leading": 0.95}], anchor="middle")
ty2 = cy + 2 * th + 8 + 6
text(RX, ty2, RW, 26, [
    {"runs": [R("Tech stack:  ", size=14, color=BLUE, bold=True), R("ESP32-CAM · Edge Impulse (TinyML) · MQTT / HiveMQ · Cloudinary · Firebase · Deno relay · React Native + Expo · TLS / WebSocket", size=13.5)], "after": 4, "leading": 1.02},
    {"runs": [R("Why it's novel:  ", size=13.5, color=BLUE, bold=True), R("Edge AI on a microcontroller (privacy + zero inference cost) + an outbound-only architecture that works on any network — the pattern commercial IoT (AWS IoT, Ring/Nest) uses, rebuilt at $0/month.", size=13, italic=True)], "leading": 1.02}])
# image row
iry = ty2 + 30
iw3 = (RW - 2 * 8) / 3
img_card(P("local_flowchart.png"), RX, iry, iw3, 120, "Detection pipeline (on-device)")
slot(RX + iw3 + 8, iry, iw3, 120, "App screenshot", "dashboard / live view", "assets/app_screenshot.png")
slot(RX + 2 * (iw3 + 8), iry, iw3, 120, "Real detection", "one clean capture", "assets/detection_sample.jpg")

# ================= FOOTER =================
fy = 800
rect(0, fy, 594, 41, fill=BLUE)
text(12, fy, 470, 41, [
    {"runs": [R("Team:  ", size=13, color=TEAL, bold=True), R("Murtaza Khalid (BSCE22004) · Huzaifa Khan (BSCE22025) · Haseeb Ahmed (BSCE22048)", size=12.5, color=WHITE)], "after": 2},
    {"runs": [R("Supervisor:  ", size=12, color=TEAL, bold=True), R("Dr. Rehan Hafiz", size=12, color=WHITE), R("     Co-Supervisor:  ", size=12, color=TEAL, bold=True), R("Dr. Rehan Ahmed", size=12, color=WHITE)], "after": 2},
    {"runs": [R("Dept. of Computer & Software Engineering, ITU Lahore   ·   SparkUp Innovation Summit 2026 (19 June 2026)   ·   UN SDG 9 · 11 · 16", size=11, color=LBLUE)]},
], anchor="middle")
qsz = 31
qx_repo = 512
rect(qx_repo - 2, fy + 4, qsz + 4, qsz + 4, fill=WHITE, rounded=True)
place_image(qr_repo, qx_repo, fy + 6, qsz, qsz)
text(qx_repo - 6, fy + qsz + 7, qsz + 12, 8, [{"runs": [R("GitHub repo", size=8.5, color=WHITE)], "align": PP_ALIGN.CENTER}])
qx_demo = 552
slot(qx_demo, fy + 4, qsz + 6, qsz + 2, "Demo QR", "add URL", "")
text(qx_demo - 4, fy + qsz + 7, qsz + 14, 8, [{"runs": [R("Demo video", size=8.5, color=WHITE)], "align": PP_ALIGN.CENTER}])

prs.save(OUT)
print("SAVED:", OUT)
